"""
Performance Metrics Collection and Analysis System
Tracks accuracy, latency, efficiency metrics for the project
"""

import time
import psutil
import json
import os
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Any, Tuple
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import numpy as np


class MetricsCollector:
    """Main class for collecting and storing performance metrics"""
    
    def __init__(self, project_name: str = "Quantum Flow"):
        self.project_name = project_name
        self.metrics = {
            "accuracy": [],
            "latency": [],
            "efficiency": [],
            "timestamp": datetime.now().isoformat()
        }
        self.metrics_file = Path("metrics_data.json")
        
    def record_accuracy(self, model_name: str, accuracy: float, precision: float, 
                       recall: float, f1_score: float, confusion_matrix: Dict = None):
        """Record model accuracy metrics"""
        accuracy_record = {
            "model_name": model_name,
            "accuracy": accuracy,
            "precision": precision,
            "recall": recall,
            "f1_score": f1_score,
            "confusion_matrix": confusion_matrix,
            "timestamp": datetime.now().isoformat()
        }
        self.metrics["accuracy"].append(accuracy_record)
        print(f"✓ Recorded accuracy for {model_name}: {accuracy:.2%}")
        return accuracy_record

    def record_latency(self, operation: str, latency_ms: float, throughput: float = None):
        """Record operation latency and throughput"""
        latency_record = {
            "operation": operation,
            "latency_ms": latency_ms,
            "throughput": throughput,  # items/sec
            "timestamp": datetime.now().isoformat()
        }
        self.metrics["latency"].append(latency_record)
        print(f"✓ Recorded latency for {operation}: {latency_ms:.2f}ms")
        return latency_record

    def record_efficiency(self, operation: str, memory_mb: float, cpu_percent: float, 
                         gpu_percent: float = None, model_size_mb: float = None):
        """Record resource efficiency metrics"""
        efficiency_record = {
            "operation": operation,
            "memory_mb": memory_mb,
            "cpu_percent": cpu_percent,
            "gpu_percent": gpu_percent,
            "model_size_mb": model_size_mb,
            "timestamp": datetime.now().isoformat()
        }
        self.metrics["efficiency"].append(efficiency_record)
        print(f"✓ Recorded efficiency for {operation}: {memory_mb:.2f}MB RAM, {cpu_percent:.1f}% CPU")
        return efficiency_record

    def measure_function_latency(self, func):
        """Decorator to measure function execution time"""
        def wrapper(*args, **kwargs):
            start = time.perf_counter()
            process = psutil.Process()
            mem_before = process.memory_info().rss / 1024 / 1024
            
            result = func(*args, **kwargs)
            
            end = time.perf_counter()
            mem_after = process.memory_info().rss / 1024 / 1024
            
            latency_ms = (end - start) * 1000
            memory_used = mem_after - mem_before
            
            self.record_latency(
                func.__name__, 
                latency_ms,
                throughput=1000/latency_ms  # items per second
            )
            
            return result
        return wrapper

    def save_metrics(self, filepath: str = None):
        """Save metrics to JSON file"""
        if filepath is None:
            filepath = str(self.metrics_file)
        
        with open(filepath, 'w') as f:
            json.dump(self.metrics, f, indent=2)
        print(f"✓ Metrics saved to {filepath}")

    def load_metrics(self, filepath: str = None):
        """Load metrics from JSON file"""
        if filepath is None:
            filepath = str(self.metrics_file)
        
        if Path(filepath).exists():
            with open(filepath, 'r') as f:
                self.metrics = json.load(f)
            print(f"✓ Metrics loaded from {filepath}")
        else:
            print(f"✗ Metrics file not found: {filepath}")

    def get_accuracy_summary(self) -> Dict:
        """Get summary statistics for accuracy metrics"""
        if not self.metrics["accuracy"]:
            return {}
        
        summary = {}
        for record in self.metrics["accuracy"]:
            model = record["model_name"]
            summary[model] = {
                "accuracy": record["accuracy"],
                "precision": record["precision"],
                "recall": record["recall"],
                "f1_score": record["f1_score"]
            }
        return summary

    def get_latency_summary(self) -> Dict:
        """Get summary statistics for latency metrics"""
        if not self.metrics["latency"]:
            return {}
        
        summary = {}
        for record in self.metrics["latency"]:
            operation = record["operation"]
            if operation not in summary:
                summary[operation] = []
            summary[operation].append(record["latency_ms"])
        
        # Calculate statistics
        stats = {}
        for operation, latencies in summary.items():
            stats[operation] = {
                "min": min(latencies),
                "max": max(latencies),
                "mean": np.mean(latencies),
                "median": np.median(latencies),
                "p95": np.percentile(latencies, 95),
                "p99": np.percentile(latencies, 99),
                "count": len(latencies)
            }
        return stats

    def get_efficiency_summary(self) -> Dict:
        """Get summary statistics for efficiency metrics"""
        if not self.metrics["efficiency"]:
            return {}
        
        summary = {}
        for record in self.metrics["efficiency"]:
            operation = record["operation"]
            summary[operation] = {
                "memory_mb": record["memory_mb"],
                "cpu_percent": record["cpu_percent"],
                "gpu_percent": record["gpu_percent"],
                "model_size_mb": record["model_size_mb"]
            }
        return summary

    def print_summary(self):
        """Print metrics summary to console"""
        print("\n" + "="*80)
        print(f"PERFORMANCE METRICS SUMMARY - {self.project_name}")
        print("="*80)
        
        # Accuracy Summary
        acc_summary = self.get_accuracy_summary()
        if acc_summary:
            print("\n📊 ACCURACY METRICS:")
            print("-" * 80)
            for model, metrics in acc_summary.items():
                print(f"  {model}:")
                print(f"    Accuracy:  {metrics['accuracy']:.2%}")
                print(f"    Precision: {metrics['precision']:.2%}")
                print(f"    Recall:    {metrics['recall']:.2%}")
                print(f"    F1-Score:  {metrics['f1_score']:.2%}")
        
        # Latency Summary
        lat_summary = self.get_latency_summary()
        if lat_summary:
            print("\n⚡ LATENCY METRICS (milliseconds):")
            print("-" * 80)
            for operation, stats in lat_summary.items():
                print(f"  {operation}:")
                print(f"    Mean:   {stats['mean']:.2f}ms | Median: {stats['median']:.2f}ms")
                print(f"    P95:    {stats['p95']:.2f}ms | P99:    {stats['p99']:.2f}ms")
                print(f"    Range:  {stats['min']:.2f}ms - {stats['max']:.2f}ms")
        
        # Efficiency Summary
        eff_summary = self.get_efficiency_summary()
        if eff_summary:
            print("\n💾 EFFICIENCY METRICS:")
            print("-" * 80)
            for operation, metrics in eff_summary.items():
                print(f"  {operation}:")
                print(f"    Memory: {metrics['memory_mb']:.2f}MB | CPU: {metrics['cpu_percent']:.1f}%")
                if metrics['gpu_percent']:
                    print(f"    GPU: {metrics['gpu_percent']:.1f}%")
                if metrics['model_size_mb']:
                    print(f"    Model Size: {metrics['model_size_mb']:.2f}MB")
        
        print("\n" + "="*80 + "\n")


class MetricsVisualizer:
    """Generate performance charts and graphs"""
    
    def __init__(self, metrics_collector: MetricsCollector):
        self.collector = metrics_collector
        self.figures = []
        sns.set_style("whitegrid")
        plt.rcParams['figure.figsize'] = (12, 6)
        
    def create_accuracy_comparison_chart(self) -> plt.Figure:
        """Create accuracy comparison bar chart"""
        acc_summary = self.collector.get_accuracy_summary()
        
        if not acc_summary:
            return None
        
        models = list(acc_summary.keys())
        accuracy = [acc_summary[m]['accuracy'] * 100 for m in models]
        precision = [acc_summary[m]['precision'] * 100 for m in models]
        recall = [acc_summary[m]['recall'] * 100 for m in models]
        f1 = [acc_summary[m]['f1_score'] * 100 for m in models]
        
        fig, ax = plt.subplots(figsize=(12, 6))
        
        x = np.arange(len(models))
        width = 0.2
        
        ax.bar(x - 1.5*width, accuracy, width, label='Accuracy', alpha=0.8)
        ax.bar(x - 0.5*width, precision, width, label='Precision', alpha=0.8)
        ax.bar(x + 0.5*width, recall, width, label='Recall', alpha=0.8)
        ax.bar(x + 1.5*width, f1, width, label='F1-Score', alpha=0.8)
        
        ax.set_xlabel('Model', fontsize=12, fontweight='bold')
        ax.set_ylabel('Score (%)', fontsize=12, fontweight='bold')
        ax.set_title('Model Accuracy Metrics Comparison', fontsize=14, fontweight='bold')
        ax.set_xticks(x)
        ax.set_xticklabels(models)
        ax.legend()
        ax.set_ylim([0, 105])
        ax.grid(axis='y', alpha=0.3)
        
        plt.tight_layout()
        self.figures.append(('accuracy_comparison', fig))
        return fig

    def create_latency_chart(self) -> plt.Figure:
        """Create latency distribution chart"""
        lat_summary = self.collector.get_latency_summary()
        
        if not lat_summary:
            return None
        
        fig, ax = plt.subplots(figsize=(12, 6))
        
        operations = list(lat_summary.keys())
        means = [lat_summary[op]['mean'] for op in operations]
        p95 = [lat_summary[op]['p95'] for op in operations]
        p99 = [lat_summary[op]['p99'] for op in operations]
        
        x = np.arange(len(operations))
        width = 0.25
        
        ax.bar(x - width, means, width, label='Mean', alpha=0.8)
        ax.bar(x, p95, width, label='P95', alpha=0.8)
        ax.bar(x + width, p99, width, label='P99', alpha=0.8)
        
        ax.set_xlabel('Operation', fontsize=12, fontweight='bold')
        ax.set_ylabel('Latency (ms)', fontsize=12, fontweight='bold')
        ax.set_title('API Latency Performance', fontsize=14, fontweight='bold')
        ax.set_xticks(x)
        ax.set_xticklabels(operations, rotation=45, ha='right')
        ax.legend()
        ax.grid(axis='y', alpha=0.3)
        
        plt.tight_layout()
        self.figures.append(('latency_chart', fig))
        return fig

    def create_throughput_chart(self) -> plt.Figure:
        """Create throughput comparison chart"""
        lat_data = self.collector.metrics.get("latency", [])
        
        if not lat_data:
            return None
        
        operations = {}
        for record in lat_data:
            op = record["operation"]
            if record.get("throughput"):
                if op not in operations:
                    operations[op] = []
                operations[op].append(record["throughput"])
        
        if not operations:
            return None
        
        fig, ax = plt.subplots(figsize=(12, 6))
        
        ops = list(operations.keys())
        throughputs = [np.mean(operations[op]) for op in ops]
        
        colors = plt.cm.viridis(np.linspace(0, 1, len(ops)))
        bars = ax.barh(ops, throughputs, color=colors, alpha=0.8)
        
        ax.set_xlabel('Throughput (items/sec)', fontsize=12, fontweight='bold')
        ax.set_title('API Throughput Performance', fontsize=14, fontweight='bold')
        ax.grid(axis='x', alpha=0.3)
        
        # Add value labels
        for i, bar in enumerate(bars):
            width = bar.get_width()
            ax.text(width, bar.get_y() + bar.get_height()/2, 
                   f'{width:.1f}', ha='left', va='center', fontweight='bold')
        
        plt.tight_layout()
        self.figures.append(('throughput_chart', fig))
        return fig

    def create_efficiency_chart(self) -> plt.Figure:
        """Create resource efficiency chart"""
        eff_summary = self.collector.get_efficiency_summary()
        
        if not eff_summary:
            return None
        
        operations = list(eff_summary.keys())
        memory = [eff_summary[op]['memory_mb'] for op in operations]
        cpu = [eff_summary[op]['cpu_percent'] for op in operations]
        
        fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(14, 5))
        
        # Memory chart
        ax1.bar(operations, memory, color='steelblue', alpha=0.8, edgecolor='navy', linewidth=1.5)
        ax1.set_ylabel('Memory (MB)', fontweight='bold', fontsize=11)
        ax1.set_title('Memory Usage per Operation', fontweight='bold', fontsize=12)
        ax1.tick_params(axis='x', rotation=45)
        ax1.grid(axis='y', alpha=0.3, linestyle='--')
        
        # CPU chart
        ax2.bar(operations, cpu, color='coral', alpha=0.8, edgecolor='darkred', linewidth=1.5)
        ax2.set_ylabel('CPU Usage (%)', fontweight='bold', fontsize=11)
        ax2.set_title('CPU Utilization per Operation', fontweight='bold', fontsize=12)
        ax2.tick_params(axis='x', rotation=45)
        ax2.grid(axis='y', alpha=0.3, linestyle='--')
        
        plt.suptitle('Resource Efficiency Metrics', fontsize=16, fontweight='bold', y=1.02)
        plt.tight_layout()
        self.figures.append(('efficiency_chart', fig))
        return fig

    def save_all_figures(self, output_dir: str = "metrics_charts"):
        """Save all generated figures as PNG"""
        Path(output_dir).mkdir(exist_ok=True)
        
        for name, fig in self.figures:
            filepath = Path(output_dir) / f"{name}.png"
            fig.savefig(filepath, dpi=300, bbox_inches='tight')
            print(f"✓ Saved chart: {filepath}")
        
        return output_dir


class PowerPointReportGenerator:
    """Generate PowerPoint presentations from metrics"""
    
    def __init__(self, metrics_collector: MetricsCollector, visualizer: MetricsVisualizer):
        self.collector = metrics_collector
        self.visualizer = visualizer
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
    def add_title_slide(self, title: str, subtitle: str = "Performance Metrics Report"):
        """Add title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Add background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(33, 66, 99)  # Dark blue
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_p = title_frame.paragraphs[0]
        title_p.font.size = Pt(54)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        title_p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.font.size = Pt(28)
        subtitle_p.font.color.rgb = RGBColor(200, 200, 200)
        subtitle_p.alignment = PP_ALIGN.CENTER
        
        # Date
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
        date_frame = date_box.text_frame
        date_frame.text = datetime.now().strftime("%B %d, %Y")
        date_p = date_frame.paragraphs[0]
        date_p.font.size = Pt(16)
        date_p.font.color.rgb = RGBColor(150, 150, 150)
        date_p.alignment = PP_ALIGN.CENTER

    def add_metrics_table_slide(self, title: str, data: Dict):
        """Add slide with metrics table"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_p = title_frame.paragraphs[0]
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        
        # Create table
        if not data:
            text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(4))
            tf = text_box.text_frame
            tf.text = "No data available"
            return
        
        # Determine rows and columns
        rows = len(data) + 1
        cols = max(len(v) for v in data.values()) + 1 if isinstance(data, dict) else 2
        
        left = Inches(0.75)
        top = Inches(1.2)
        width = Inches(8.5)
        height = Inches(5)
        
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Header
        col_idx = 0
        for col_val in list(data.keys()):
            cell = table_shape.cell(0, col_idx)
            cell.text = col_val
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(33, 66, 99)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)
                    run.font.bold = True
            col_idx += 1
        
        # Data rows
        row_idx = 1
        for key, values in data.items():
            if isinstance(values, dict):
                table_shape.cell(row_idx, 0).text = key
                col_idx = 1
                for val in values.values():
                    table_shape.cell(row_idx, col_idx).text = f"{val:.2f}" if isinstance(val, float) else str(val)
                    col_idx += 1
                row_idx += 1

    def add_chart_slide(self, title: str, chart_path: str):
        """Add slide with chart image"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_p = title_frame.paragraphs[0]
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        
        # Add image
        try:
            left = Inches(0.5)
            top = Inches(1.0)
            height = Inches(6)
            slide.shapes.add_picture(chart_path, left, top, height=height)
        except Exception as e:
            print(f"Error adding image: {e}")

    def add_summary_slide(self):
        """Add summary statistics slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = "Performance Summary"
        title_p = title_frame.paragraphs[0]
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        
        # Summary text
        summary_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(8.5), Inches(5.5))
        tf = summary_box.text_frame
        tf.word_wrap = True
        
        # Accuracy summary
        acc_summary = self.collector.get_accuracy_summary()
        if acc_summary:
            p = tf.add_paragraph()
            p.text = "📊 Accuracy Metrics"
            p.font.size = Pt(18)
            p.font.bold = True
            p.level = 0
            
            for model, metrics in acc_summary.items():
                p = tf.add_paragraph()
                p.text = f"{model}: {metrics['accuracy']:.2%} accuracy"
                p.font.size = Pt(14)
                p.level = 1
        
        # Latency summary
        lat_summary = self.collector.get_latency_summary()
        if lat_summary:
            p = tf.add_paragraph()
            p.text = "⚡ Latency Performance"
            p.font.size = Pt(18)
            p.font.bold = True
            p.level = 0
            p.space_before = Pt(12)
            
            for operation, stats in lat_summary.items():
                p = tf.add_paragraph()
                p.text = f"{operation}: {stats['mean']:.2f}ms (mean)"
                p.font.size = Pt(14)
                p.level = 1

    def generate_report(self, output_file: str = "Performance_Report.pptx"):
        """Generate complete PowerPoint report"""
        self.add_title_slide(self.collector.project_name, "Performance Metrics Report")
        
        # Add metric charts
        self.visualizer.create_accuracy_comparison_chart()
        self.visualizer.create_latency_chart()
        self.visualizer.create_throughput_chart()
        self.visualizer.create_efficiency_chart()
        
        # Save charts
        chart_dir = self.visualizer.save_all_figures()
        
        # Add slides
        self.add_chart_slide("Accuracy Comparison", Path(chart_dir) / "accuracy_comparison.png")
        self.add_chart_slide("Latency Performance", Path(chart_dir) / "latency_chart.png")
        self.add_chart_slide("Throughput Analysis", Path(chart_dir) / "throughput_chart.png")
        self.add_chart_slide("Efficiency Metrics", Path(chart_dir) / "efficiency_chart.png")
        self.add_summary_slide()
        
        # Add metrics tables
        acc_summary = self.collector.get_accuracy_summary()
        if acc_summary:
            self.add_metrics_table_slide("Accuracy Metrics Table", acc_summary)
        
        # Save presentation
        self.prs.save(output_file)
        print(f"✓ PowerPoint report saved: {output_file}")
        return output_file


# Example usage function
def example_usage():
    """Example of how to use the metrics system"""
    
    # Initialize collector
    collector = MetricsCollector("Quantum Flow - Healthcare AI")
    
    # Record accuracy metrics (simulated)
    collector.record_accuracy(
        model_name="Baseline CNN",
        accuracy=0.852,
        precision=0.841,
        recall=0.863,
        f1_score=0.852,
        confusion_matrix={"TP": 45, "FP": 8, "TN": 92, "FN": 5}
    )
    
    collector.record_accuracy(
        model_name="Quantum Enhanced",
        accuracy=0.893,
        precision=0.889,
        recall=0.897,
        f1_score=0.893,
        confusion_matrix={"TP": 47, "FP": 6, "TN": 94, "FN": 3}
    )
    
    # Record latency metrics (simulated)
    for i in range(10):
        collector.record_latency("Image Upload", 45 + np.random.normal(0, 5), throughput=22)
        collector.record_latency("Vector Search", 35 + np.random.normal(0, 4), throughput=28)
        collector.record_latency("Inference", 120 + np.random.normal(0, 10), throughput=8)
    
    # Record efficiency metrics
    collector.record_efficiency(
        operation="Image Preprocessing",
        memory_mb=256.5,
        cpu_percent=45.2,
        gpu_percent=62.1,
        model_size_mb=128.3
    )
    
    collector.record_efficiency(
        operation="Vector Embedding",
        memory_mb=512.0,
        cpu_percent=38.5,
        gpu_percent=78.3,
        model_size_mb=256.0
    )
    
    # Print summary
    collector.print_summary()
    
    # Save metrics
    collector.save_metrics("metrics_data.json")
    
    # Create visualizations
    visualizer = MetricsVisualizer(collector)
    visualizer.create_accuracy_comparison_chart()
    visualizer.create_latency_chart()
    visualizer.create_throughput_chart()
    visualizer.create_efficiency_chart()
    
    # Generate PowerPoint report
    report_gen = PowerPointReportGenerator(collector, visualizer)
    report_gen.generate_report("Performance_Metrics_Report.pptx")
    
    print("\n✓ Complete metrics collection and reporting system initialized!")
    print("  - Charts saved in 'metrics_charts/' directory")
    print("  - PowerPoint report: Performance_Metrics_Report.pptx")
    print("  - JSON data: metrics_data.json")


if __name__ == "__main__":
    example_usage()
